#!/usr/bin/env python
# coding: utf-8

# In[21]:


import selenium
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.keys import Keys
import warnings, time
warnings.filterwarnings('ignore')

# In[22]:

s=Service('C:/Users/chromedriver.exe')
driver = webdriver.Chrome(service=s)
driver.get("https://topic.udn.com/event/COVID19_Taiwan")
pause_time = 0.5

# Get scroll height
for i in range(20):
    driver.execute_script(f"window.scrollTo(0, {200*i})") 
    time.sleep(pause_time)


# In[23]:
driver.switch_to.frame(2)

# In[24]:

from selenium.webdriver.support.ui import WebDriverWait

try:
    imgs = WebDriverWait(driver, timeout=50).until(lambda d: d.find_elements_by_class_name("ItemContent-image"))
    imageUrl = imgs[1].get_attribute("src")
except:
    driver.refresh


# In[25]:


import re
stat_num = ''
for element in driver.find_elements_by_class_name("__ig-alignLeft"):
    text = element.text.replace('  ',' ').strip()
    text = re.sub(" +", " ",text)
    text = re.sub("[^0-9.%, ]", "",text).strip()
    if len(text.split(' ')) == 3:
        stat_num = text
        
stat_num = stat_num.replace(" ", "              ")
stat_num


# In[26]:


import requests

img_data = requests.get(imageUrl).content
with open('image_name.jpg', 'wb') as handler:
    handler.write(img_data)


# In[27]:


from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# In[28]:


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# create title textbox
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(0.2), Inches(8), Inches(1)
)

fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(44,177,153)

shape.line.fill.background()
title = shape.text_frame

pic = slide.shapes.add_picture('image_name.jpg', Inches(3.5), Inches(1.4), Inches(9), Inches(6.1))

# create subtitle textbox
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), Inches(7.7), Inches(9.5), Inches(1.1)
)

fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(160,155,188)

shape.line.fill.background()
subtitle = shape.text_frame

month = datetime.now().strftime("%m")
date = datetime.now().strftime("%d")

if int(month) < 10:
    month = month[1:]

if int(date) < 10:
    date = date[1:]

title.text = month + "月" + date + "日" +"台灣疫情數據"
subtitle.text = '累計確診           累計死亡           死亡率\n' + stat_num

title.paragraphs[0].font.size = Pt(26)
title.paragraphs[0].font.name = '微軟正黑體'
title.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle.paragraphs[0].font.size = Pt(26)
subtitle.paragraphs[0].font.name = '微軟正黑體'
subtitle.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle.paragraphs[1].font.size = Pt(26)
subtitle.paragraphs[1].font.name = '微軟正黑體'
subtitle.paragraphs[1].alignment = PP_ALIGN.CENTER

# prs.save('test.pptx')


# In[29]:

# driver = webdriver.Chrome(executable_path='C:/Users/chromedriver.exe')
# driver.get("https://covid-19.nchc.org.tw/")
# pause_time = 0.5

# element = WebDriverWait(driver, 3000).until(
#     EC.presence_of_element_located((By.XPATH, '/html/body/div[2]/div/div[3]/p/span/small'))
# )

# cases = element.text.split(" ")[1]

# covid_data = f'{month}/{date},{cases}'

# cases

# In[30]:


# import pandas as pd
# df = pd.read_excel('stat.xlsx')
# df = pd.read_csv('covid_case.csv')


# In[31]:


# with open('covid_case.csv', mode='a') as file:
#     file.writelines(covid_data + '\n')

# In[32]:


import pandas as pd

file = open('covid_case.csv')
df = pd.read_csv(file)

# import gspread
# from oauth2client.service_account import ServiceAccountCredentials

# #credentials to the account
# cred = ServiceAccountCredentials.from_json_keyfile_name(r'C:\Users\u0198\Desktop\AutoDev\udn_scrap\cred.json') ;
# # authorize the clientsheet 
# client = gspread.authorize(cred)

# sh = client.open('coviddata')
# worksheet = sh.worksheet('2022')

# import pandas as pd
# df = pd.DataFrame(worksheet.get_all_records())


# In[33]:


from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

slide2 = prs.slides.add_slide(blank_slide_layout)

shape = slide2.shapes.add_textbox(Inches(1.8), Inches(0.4), Inches(6), Inches(1))
title = shape.text_frame

title.text = "2022、2023"
title.paragraphs[0].font.bold = True
title.paragraphs[0].font.size = Pt(32)
title.paragraphs[0].font.name = '微軟正黑體'
title.paragraphs[0].font.color.rgb = RGBColor(163,181,164)

shape = slide2.shapes.add_textbox(Inches(4.5), Inches(0.4), Inches(6), Inches(1))
title2 = shape.text_frame

title2.text = "每日新增本土確診數"
title2.paragraphs[0].font.bold = True
title2.paragraphs[0].font.size = Pt(32)
title2.paragraphs[0].font.name = '微軟正黑體'

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = df['date']
chart_data.add_series( "" , df['value'] )

# add chart to slide --------------------
x, y, cx, cy = Inches(1.8), Inches(0.7), Inches(12.5), Inches(5.8)
graphic_frame = slide2.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

chart = graphic_frame.chart


# In[34]:


fill = chart.plots[0].series[0].format.fill #fill the legend as well
fill.solid()
fill.fore_color.rgb = RGBColor(163,181,164)

chart.plots[0].gap_width = 50
points = chart.plots[0].series[0].points
for point in points:
    fill = point.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(163,181,164)


# In[35]:


# prs = Presentation()
# prs.slide_width = Inches(16)
# prs.slide_height = Inches(9)
# blank_slide_layout = prs.slide_layouts[6]
# slide = prs.slides.add_slide(blank_slide_layout)


# In[36]:

shape = slide2.shapes.add_textbox(Inches(1.8), Inches(6.6), Inches(12), Inches(1))
title = shape.text_frame
title.text = f'2/20起起室內可免戴口罩，開放北辦及廠區不戴口罩。'
# title.paragraphs[0].font.color.rgb = RGBColor(0,0,255);

for p in title.paragraphs:
  p.font.size = Pt(26)
  p.font.name = "微軟正黑體"

shape = slide2.shapes.add_textbox(Inches(1.8), Inches(7.1), Inches(7), Inches(1))
title = shape.text_frame

title.text = '在外用餐，與北辦用餐都要注意個人衛生管理。'
title.paragraphs[0].font.size = Pt(26)
title.paragraphs[0].font.name = "微軟正黑體"

# title.paragraphs[0].font.name = "微軟正黑體"
# # title.paragraphs[0].font.color.rgb = RGBColor(50,63,255)

# In[36.1]

slide3 = prs.slides.add_slide(blank_slide_layout)
shape = slide3.shapes.add_textbox(Inches(4), Inches(0.4), Inches(8), Inches(1))

title = shape.text_frame

title.text = "台灣榮成目前確診/居隔統計"
title.paragraphs[0].font.size = Pt(36)
title.paragraphs[0].font.name = '微軟正黑體'
title.paragraphs[0].alignment = PP_ALIGN.CENTER

# In[37]:

m = datetime.now().strftime("%m")
d = datetime.now().strftime("%d")

prs.save(f'新冠即時訊息({m}{d}).pptx')

# In[38]:


# from openpyxl import Workbook
# from openpyxl.utils import get_column_letter

# wb = Workbook()
# dest_filename = 'empty_book.xlsx'
# ws1 = wb.active
# ws1.title = "range names"
# ws1.cell(column=1, row=1, value="row")
# wb.save(filename = dest_filename)


# In[39]:


# create textbox
# title_shape = slide.shapes.add_shape(
#     MSO_SHAPE.RECTANGLE, Inches(5), Inches(0.4), Inches(6), Inches(0.6)
# )

# fill = title_shape.fill
# fill.solid()
# fill.fore_color.rgb = RGBColor(44,177,153)
# fill.fore_color.brightness = -0.3

# title_shape.line.fill.background()

# title = title_shape.text_frame
# title.text = '新冠即時訊息'

# title.paragraphs[0].font.size = Pt(30)
# title.paragraphs[0].font.name = 'msjh'
# title.paragraphs[0].alignment = PP_ALIGN.CENTER

